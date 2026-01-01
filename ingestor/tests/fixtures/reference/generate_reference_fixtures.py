"""Generate reference test fixtures with known expected outputs.

These fixtures are handcrafted to test complex scenarios and edge cases.
Each input file has a corresponding .expected.md file for exact comparison.

Run: python -m tests.fixtures.reference.generate_reference_fixtures
"""

import json
from pathlib import Path

REFERENCE_DIR = Path(__file__).parent


def generate_complex_docx():
    """Generate complex DOCX with nested tables, images, headers, etc."""
    try:
        from docx import Document
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Inches, Pt, RGBColor

        doc = Document()

        # Document title
        title = doc.add_heading("Complex Document Test", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Introduction paragraph with formatting
        intro = doc.add_paragraph()
        intro.add_run("This document tests ").bold = False
        intro.add_run("complex formatting").bold = True
        intro.add_run(" including ")
        run = intro.add_run("italic")
        run.italic = True
        intro.add_run(" and ")
        run = intro.add_run("bold italic")
        run.bold = True
        run.italic = True
        intro.add_run(" text.")

        # Section 1: Nested Tables
        doc.add_heading("Section 1: Tables", level=1)
        doc.add_paragraph("This section contains simple and nested tables.")

        # Simple table
        doc.add_heading("Simple Table", level=2)
        table1 = doc.add_table(rows=4, cols=3)
        table1.style = 'Table Grid'
        headers = ["Name", "Age", "City"]
        for i, header in enumerate(headers):
            table1.cell(0, i).text = header
        data = [
            ["Alice", "30", "New York"],
            ["Bob", "25", "Los Angeles"],
            ["Charlie", "35", "Chicago"],
        ]
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, cell_data in enumerate(row_data):
                table1.cell(row_idx, col_idx).text = cell_data

        doc.add_paragraph()  # Spacing

        # Nested table (table within table)
        doc.add_heading("Nested Table", level=2)
        outer_table = doc.add_table(rows=2, cols=2)
        outer_table.style = 'Table Grid'
        outer_table.cell(0, 0).text = "Cell A1"
        outer_table.cell(0, 1).text = "Cell B1"
        outer_table.cell(1, 0).text = "Cell A2 with content"

        # Add inner table to cell B2
        inner_cell = outer_table.cell(1, 1)
        inner_para = inner_cell.paragraphs[0]
        inner_para.add_run("Inner table below:")
        inner_table = inner_cell.add_table(rows=2, cols=2)
        inner_table.cell(0, 0).text = "X1"
        inner_table.cell(0, 1).text = "Y1"
        inner_table.cell(1, 0).text = "X2"
        inner_table.cell(1, 1).text = "Y2"

        # Section 2: Lists
        doc.add_heading("Section 2: Lists", level=1)

        doc.add_heading("Bullet List", level=2)
        for item in ["First item", "Second item", "Third item"]:
            doc.add_paragraph(item, style='List Bullet')

        doc.add_heading("Numbered List", level=2)
        for item in ["Step one", "Step two", "Step three"]:
            doc.add_paragraph(item, style='List Number')

        # Section 3: Code and Special Text
        doc.add_heading("Section 3: Special Content", level=1)

        doc.add_heading("Code Block", level=2)
        code_para = doc.add_paragraph()
        code_run = code_para.add_run("def hello():\n    print('Hello, World!')")
        code_run.font.name = 'Courier New'
        code_run.font.size = Pt(10)

        doc.add_heading("Hyperlinks", level=2)
        link_para = doc.add_paragraph("Visit ")
        # Note: python-docx doesn't have easy hyperlink support, just add text
        link_para.add_run("https://example.com").underline = True
        link_para.add_run(" for more info.")

        # Section 4: Multiple Headings Levels
        doc.add_heading("Section 4: Heading Hierarchy", level=1)
        doc.add_heading("Subsection 4.1", level=2)
        doc.add_paragraph("Content under 4.1")
        doc.add_heading("Sub-subsection 4.1.1", level=3)
        doc.add_paragraph("Content under 4.1.1")
        doc.add_heading("Subsection 4.2", level=2)
        doc.add_paragraph("Content under 4.2")

        # Save
        doc.save(REFERENCE_DIR / "complex.docx")
        print("Created complex.docx")

        # Create expected output
        expected_md = """# Complex Document Test

This document tests **complex formatting** including *italic* and ***bold italic*** text.

# Section 1: Tables

This section contains simple and nested tables.

## Simple Table

| Name | Age | City |
|---|---|---|
| Alice | 30 | New York |
| Bob | 25 | Los Angeles |
| Charlie | 35 | Chicago |

## Nested Table

| Cell A1 | Cell B1 |
|---|---|
| Cell A2 with content | Inner table below: |

| X1 | Y1 |
|---|---|
| X2 | Y2 |

# Section 2: Lists

## Bullet List

- First item
- Second item
- Third item

## Numbered List

1. Step one
2. Step two
3. Step three

# Section 3: Special Content

## Code Block

```
def hello():
    print('Hello, World!')
```

## Hyperlinks

Visit https://example.com for more info.

# Section 4: Heading Hierarchy

## Subsection 4.1

Content under 4.1

### Sub-subsection 4.1.1

Content under 4.1.1

## Subsection 4.2

Content under 4.2
"""
        (REFERENCE_DIR / "complex.docx.expected.md").write_text(expected_md.strip())
        print("Created complex.docx.expected.md")

    except ImportError as e:
        print(f"python-docx not installed, skipping DOCX generation: {e}")


def generate_complex_xlsx():
    """Generate complex XLSX with merged cells, formulas, multiple sheets."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

        wb = Workbook()

        # Sheet 1: Sales Data with formulas
        ws1 = wb.active
        ws1.title = "Sales"

        # Headers
        ws1['A1'] = "Product"
        ws1['B1'] = "Q1"
        ws1['C1'] = "Q2"
        ws1['D1'] = "Q3"
        ws1['E1'] = "Q4"
        ws1['F1'] = "Total"

        # Style headers
        for col in 'ABCDEF':
            ws1[f'{col}1'].font = Font(bold=True)

        # Data with formulas
        products = [
            ("Widget", 100, 150, 200, 180),
            ("Gadget", 80, 120, 90, 110),
            ("Gizmo", 200, 180, 220, 250),
        ]

        for i, (product, q1, q2, q3, q4) in enumerate(products, start=2):
            ws1[f'A{i}'] = product
            ws1[f'B{i}'] = q1
            ws1[f'C{i}'] = q2
            ws1[f'D{i}'] = q3
            ws1[f'E{i}'] = q4
            ws1[f'F{i}'] = f'=SUM(B{i}:E{i})'  # Formula

        # Total row
        ws1['A5'] = "TOTAL"
        ws1['A5'].font = Font(bold=True)
        for _col_idx, col in enumerate('BCDEF', start=2):
            ws1[f'{col}5'] = f'=SUM({col}2:{col}4)'

        # Sheet 2: Merged Cells
        ws2 = wb.create_sheet("Summary")

        # Merge title
        ws2.merge_cells('A1:D1')
        ws2['A1'] = "Quarterly Sales Summary"
        ws2['A1'].font = Font(bold=True, size=14)
        ws2['A1'].alignment = Alignment(horizontal='center')

        # Merge category header
        ws2.merge_cells('A3:B3')
        ws2['A3'] = "Category"
        ws2.merge_cells('C3:D3')
        ws2['C3'] = "Performance"

        ws2['A4'] = "Electronics"
        ws2['B4'] = "High"
        ws2['C4'] = "Good"
        ws2['D4'] = "+15%"

        ws2['A5'] = "Software"
        ws2['B5'] = "Medium"
        ws2['C5'] = "Excellent"
        ws2['D5'] = "+25%"

        # Sheet 3: Various Data Types
        ws3 = wb.create_sheet("Data Types")
        ws3['A1'] = "Type"
        ws3['B1'] = "Value"

        ws3['A2'] = "Integer"
        ws3['B2'] = 42

        ws3['A3'] = "Float"
        ws3['B3'] = 3.14159

        ws3['A4'] = "Percentage"
        ws3['B4'] = 0.75
        ws3['B4'].number_format = '0%'

        ws3['A5'] = "Currency"
        ws3['B5'] = 1234.56
        ws3['B5'].number_format = '$#,##0.00'

        ws3['A6'] = "Date"
        from datetime import date
        ws3['B6'] = date(2024, 1, 15)

        ws3['A7'] = "Boolean"
        ws3['B7'] = True

        wb.save(REFERENCE_DIR / "complex.xlsx")
        print("Created complex.xlsx")

        # Create expected output
        expected_md = """# complex.xlsx

## Sheet: Sales

| Product | Q1 | Q2 | Q3 | Q4 | Total |
|---|---|---|---|---|---|
| Widget | 100 | 150 | 200 | 180 | 630 |
| Gadget | 80 | 120 | 90 | 110 | 400 |
| Gizmo | 200 | 180 | 220 | 250 | 850 |
| TOTAL | 380 | 450 | 510 | 540 | 1880 |

## Sheet: Summary

| Quarterly Sales Summary |  |  |  |
|---|---|---|---|
|  |  |  |  |
| Category |  | Performance |  |
| Electronics | High | Good | +15% |
| Software | Medium | Excellent | +25% |

## Sheet: Data Types

| Type | Value |
|---|---|
| Integer | 42 |
| Float | 3.14159 |
| Percentage | 75% |
| Currency | $1,234.56 |
| Date | 2024-01-15 |
| Boolean | TRUE |
"""
        (REFERENCE_DIR / "complex.xlsx.expected.md").write_text(expected_md.strip())
        print("Created complex.xlsx.expected.md")

    except ImportError as e:
        print(f"openpyxl not installed, skipping XLSX generation: {e}")


def generate_complex_pptx():
    """Generate complex PPTX with multiple slides, layouts, notes."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        prs = Presentation()

        # Slide 1: Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide1 = prs.slides.add_slide(slide_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Complex Presentation Test"
        subtitle.text = "Testing PPTX extraction capabilities"

        # Add speaker notes to slide 1
        notes_slide = slide1.notes_slide
        notes_slide.notes_text_frame.text = "Welcome to the presentation. This tests complex PPTX features."

        # Slide 2: Content with bullet points
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Bullet Points"

        body = slide2.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = "Main point one"

        p = tf.add_paragraph()
        p.text = "Main point two"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Sub-point under two"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Main point three"
        p.level = 0

        # Speaker notes for slide 2
        notes_slide = slide2.notes_slide
        notes_slide.notes_text_frame.text = "Discuss each bullet point in detail."

        # Slide 3: Two columns
        slide_layout = prs.slide_layouts[3]  # Two content
        slide3 = prs.slides.add_slide(slide_layout)
        slide3.shapes.title.text = "Two Column Layout"

        # Left column
        left = slide3.shapes.placeholders[1]
        left.text_frame.text = "Left Column"
        left.text_frame.add_paragraph().text = "Item A"
        left.text_frame.add_paragraph().text = "Item B"

        # Right column
        right = slide3.shapes.placeholders[2]
        right.text_frame.text = "Right Column"
        right.text_frame.add_paragraph().text = "Item X"
        right.text_frame.add_paragraph().text = "Item Y"

        # Slide 4: Table
        slide_layout = prs.slide_layouts[5]  # Blank
        slide4 = prs.slides.add_slide(slide_layout)

        # Add title manually
        title_shape = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.paragraphs[0].text = "Data Table"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Add table
        rows, cols = 4, 3
        table = slide4.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(2)).table

        # Headers
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Role"
        table.cell(0, 2).text = "Status"

        # Data
        data = [
            ("Alice", "Developer", "Active"),
            ("Bob", "Designer", "Active"),
            ("Charlie", "Manager", "Away"),
        ]
        for i, (name, role, status) in enumerate(data, start=1):
            table.cell(i, 0).text = name
            table.cell(i, 1).text = role
            table.cell(i, 2).text = status

        # Slide 5: Shapes and text
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])

        # Title
        title_shape = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.paragraphs[0].text = "Shapes and Text"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Add a rectangle with text
        shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1.5))
        shape.text = "Rectangle with text inside"

        # Add an oval with text
        shape2 = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2), Inches(3), Inches(1.5))
        shape2.text = "Oval shape"

        # Additional text box
        textbox = slide5.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        textbox.text_frame.text = "This is a standalone text box with important information."

        prs.save(REFERENCE_DIR / "complex.pptx")
        print("Created complex.pptx")

        # Create expected output
        expected_md = """# Complex Presentation Test

Testing PPTX extraction capabilities

---

## Slide 1: Complex Presentation Test

Testing PPTX extraction capabilities

**Speaker Notes:**
Welcome to the presentation. This tests complex PPTX features.

---

## Slide 2: Bullet Points

- Main point one
- Main point two
  - Sub-point under two
- Main point three

**Speaker Notes:**
Discuss each bullet point in detail.

---

## Slide 3: Two Column Layout

### Left Column
- Item A
- Item B

### Right Column
- Item X
- Item Y

---

## Slide 4: Data Table

| Name | Role | Status |
|---|---|---|
| Alice | Developer | Active |
| Bob | Designer | Active |
| Charlie | Manager | Away |

---

## Slide 5: Shapes and Text

Rectangle with text inside

Oval shape

This is a standalone text box with important information.
"""
        (REFERENCE_DIR / "complex.pptx.expected.md").write_text(expected_md.strip())
        print("Created complex.pptx.expected.md")

    except ImportError as e:
        print(f"python-pptx not installed, skipping PPTX generation: {e}")


def generate_edge_case_files():
    """Generate edge case test files."""

    # Empty files
    (REFERENCE_DIR / "empty.txt").write_text("")
    (REFERENCE_DIR / "empty.txt.expected.md").write_text("")
    print("Created empty.txt")

    # Whitespace only
    (REFERENCE_DIR / "whitespace_only.txt").write_text("   \n\n   \t\t\n   ")
    (REFERENCE_DIR / "whitespace_only.txt.expected.md").write_text("")
    print("Created whitespace_only.txt")

    # Unicode filename content
    unicode_content = """# Unicode Test

こんにちは (Japanese)
مرحبا (Arabic)
Привет (Russian)
🎉 Emoji support 🚀
"""
    (REFERENCE_DIR / "unicode_content.txt").write_text(unicode_content, encoding="utf-8")
    (REFERENCE_DIR / "unicode_content.txt.expected.md").write_text(unicode_content.strip())
    print("Created unicode_content.txt")

    # Very long line
    long_line = "A" * 10000
    (REFERENCE_DIR / "long_line.txt").write_text(f"# Long Line Test\n\n{long_line}\n")
    (REFERENCE_DIR / "long_line.txt.expected.md").write_text(f"# Long Line Test\n\n{long_line}")
    print("Created long_line.txt")

    # Special characters
    special_content = """# Special Characters

Quotes: "double" and 'single'
Ampersand: &
Less/Greater: < >
Backslash: \\
Pipe: |
Asterisks: * ** ***
Underscores: _ __ ___
Backticks: ` `` ```
"""
    (REFERENCE_DIR / "special_chars.txt").write_text(special_content)
    (REFERENCE_DIR / "special_chars.txt.expected.md").write_text(special_content.strip())
    print("Created special_chars.txt")

    # JSON edge cases
    json_edge = {
        "empty_string": "",
        "null_value": None,
        "nested": {
            "deep": {
                "deeper": {
                    "deepest": "value"
                }
            }
        },
        "array_of_arrays": [[1, 2], [3, 4], [5, 6]],
        "special_chars": "Line1\nLine2\tTabbed",
        "unicode": "日本語テスト",
        "numbers": {
            "int": 42,
            "float": 3.14159,
            "negative": -100,
            "scientific": 1.23e10
        },
        "boolean": {
            "true": True,
            "false": False
        }
    }
    (REFERENCE_DIR / "edge_cases.json").write_text(json.dumps(json_edge, indent=2, ensure_ascii=False))
    print("Created edge_cases.json")

    expected_json_md = """# edge_cases.json

```json
{
  "empty_string": "",
  "null_value": null,
  "nested": {
    "deep": {
      "deeper": {
        "deepest": "value"
      }
    }
  },
  "array_of_arrays": [
    [1, 2],
    [3, 4],
    [5, 6]
  ],
  "special_chars": "Line1\\nLine2\\tTabbed",
  "unicode": "日本語テスト",
  "numbers": {
    "int": 42,
    "float": 3.14159,
    "negative": -100,
    "scientific": 12300000000.0
  },
  "boolean": {
    "true": true,
    "false": false
  }
}
```
"""
    (REFERENCE_DIR / "edge_cases.json.expected.md").write_text(expected_json_md.strip())
    print("Created edge_cases.json.expected.md")

    # CSV edge cases
    csv_edge = '''name,description,value
"Simple","Plain text",100
"With, Comma","Text with, comma",200
"With ""Quotes""","Text with ""embedded"" quotes",300
"With
Newline","Multi
line
text",400
"Empty","",0
'''
    (REFERENCE_DIR / "edge_cases.csv").write_text(csv_edge)
    print("Created edge_cases.csv")

    expected_csv_md = """# edge_cases.csv

| name | description | value |
|---|---|---|
| Simple | Plain text | 100 |
| With, Comma | Text with, comma | 200 |
| With "Quotes" | Text with "embedded" quotes | 300 |
| With Newline | Multi line text | 400 |
| Empty |  | 0 |
"""
    (REFERENCE_DIR / "edge_cases.csv.expected.md").write_text(expected_csv_md.strip())
    print("Created edge_cases.csv.expected.md")


def generate_simple_text_reference():
    """Generate simple text reference for exact matching."""

    content = """# Simple Test Document

This is a paragraph with **bold** and *italic* text.

## Section One

- Bullet one
- Bullet two
- Bullet three

## Section Two

1. Numbered one
2. Numbered two
3. Numbered three

### Subsection

Some final content.
"""
    (REFERENCE_DIR / "simple.md").write_text(content.strip())
    (REFERENCE_DIR / "simple.md.expected.md").write_text(content.strip())
    print("Created simple.md reference")


def main():
    """Generate all reference fixtures."""
    print("=" * 60)
    print("Generating Reference Test Fixtures")
    print("=" * 60)
    print(f"Output directory: {REFERENCE_DIR}")
    print()

    generate_simple_text_reference()
    generate_edge_case_files()
    generate_complex_docx()
    generate_complex_xlsx()
    generate_complex_pptx()

    print()
    print("=" * 60)
    print("Done! Reference fixtures generated.")
    print("=" * 60)


if __name__ == "__main__":
    main()
