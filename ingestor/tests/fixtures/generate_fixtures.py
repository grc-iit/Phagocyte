"""Generate test fixture files.

Run this script to create all test fixtures:
    python -m tests.fixtures.generate_fixtures

Fixtures are created with known content so tests can verify
exact markdown output.
"""

import json
import struct
import wave
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def generate_text_files():
    """Generate text file fixtures."""
    documents_dir = FIXTURES_DIR / "documents"
    documents_dir.mkdir(exist_ok=True)

    # UTF-8 text file
    content = """# Test Document

This is a test document with known content.

## Section 1

Some text in section 1.

## Section 2

Some text in section 2.
"""
    (documents_dir / "sample.txt").write_text(content, encoding="utf-8")
    print("Created sample.txt")

    # Markdown file
    md_content = """# Markdown Document

This is a **markdown** document with _formatting_.

## Features

- Bullet points
- **Bold text**
- *Italic text*

## Code

```python
print("Hello, World!")
```

[Link example](https://example.com)
"""
    (documents_dir / "sample.md").write_text(md_content, encoding="utf-8")
    print("Created sample.md")

    # reStructuredText file
    rst_content = """===============
RST Document
===============

This is a reStructuredText document.

Section 1
---------

Some content in section 1.

Section 2
---------

- Item one
- Item two

``code example``
"""
    (documents_dir / "sample.rst").write_text(rst_content, encoding="utf-8")
    print("Created sample.rst")

    # Latin-1 encoded file with special characters
    latin1_content = "Héllo Wörld! Çafé résumé naïve"
    (documents_dir / "sample_latin1.txt").write_bytes(
        latin1_content.encode("latin-1")
    )
    print("Created sample_latin1.txt")


def generate_json_files():
    """Generate JSON file fixtures."""
    data_dir = FIXTURES_DIR / "data"
    data_dir.mkdir(exist_ok=True)

    # Object JSON
    obj_data = {
        "title": "Test Document",
        "version": 1,
        "items": [
            {"id": 1, "name": "Item 1"},
            {"id": 2, "name": "Item 2"},
        ],
    }
    (data_dir / "sample.json").write_text(
        json.dumps(obj_data, indent=2), encoding="utf-8"
    )
    print("Created sample.json")

    # Array JSON (tabular)
    array_data = [
        {"id": 1, "name": "Alice", "age": 30},
        {"id": 2, "name": "Bob", "age": 25},
        {"id": 3, "name": "Charlie", "age": 35},
    ]
    (data_dir / "sample_array.json").write_text(
        json.dumps(array_data, indent=2), encoding="utf-8"
    )
    print("Created sample_array.json")


def generate_xml_files():
    """Generate XML file fixtures."""
    data_dir = FIXTURES_DIR / "data"
    data_dir.mkdir(exist_ok=True)

    xml_content = """<?xml version="1.0" encoding="UTF-8"?>
<root>
    <title>Test Document</title>
    <items>
        <item id="1">First Item</item>
        <item id="2">Second Item</item>
    </items>
</root>
"""
    (data_dir / "sample.xml").write_text(xml_content, encoding="utf-8")
    print("Created sample.xml")


def generate_csv_files():
    """Generate CSV file fixtures."""
    spreadsheets_dir = FIXTURES_DIR / "spreadsheets"
    spreadsheets_dir.mkdir(exist_ok=True)

    csv_content = """name,age,city
Alice,30,New York
Bob,25,San Francisco
Charlie,35,Chicago
"""
    (spreadsheets_dir / "sample.csv").write_text(csv_content, encoding="utf-8")
    print("Created sample.csv")

    # TSV (tab-separated)
    tsv_content = "name\tage\tcity\nAlice\t30\tNew York\nBob\t25\tSan Francisco\n"
    (spreadsheets_dir / "sample.tsv").write_text(tsv_content, encoding="utf-8")
    print("Created sample.tsv")


def generate_image_files():
    """Generate image file fixtures using Pillow."""
    images_dir = FIXTURES_DIR / "images"
    images_dir.mkdir(exist_ok=True)

    try:
        from PIL import Image

        # PNG - solid red with transparency
        img_png = Image.new("RGBA", (100, 100), (255, 0, 0, 200))
        img_png.save(images_dir / "sample.png")
        print("Created sample.png")

        # JPEG - solid blue (no transparency)
        img_jpg = Image.new("RGB", (100, 100), (0, 0, 255))
        img_jpg.save(images_dir / "sample.jpg", quality=90)
        print("Created sample.jpg")

        # Image with EXIF-like properties (just a larger image)
        img_exif = Image.new("RGB", (800, 600), (0, 255, 0))
        img_exif.save(images_dir / "sample_large.jpg", quality=85)
        print("Created sample_large.jpg")

    except ImportError:
        print("Pillow not installed, skipping image generation")


def generate_xlsx_files():
    """Generate XLSX file fixtures using openpyxl."""
    spreadsheets_dir = FIXTURES_DIR / "spreadsheets"
    spreadsheets_dir.mkdir(exist_ok=True)

    try:
        from openpyxl import Workbook

        wb = Workbook()

        # Sheet 1
        ws1 = wb.active
        ws1.title = "People"
        ws1.append(["Name", "Age", "City"])
        ws1.append(["Alice", 30, "New York"])
        ws1.append(["Bob", 25, "San Francisco"])
        ws1.append(["Charlie", 35, "Chicago"])

        # Sheet 2
        ws2 = wb.create_sheet("Products")
        ws2.append(["ID", "Product", "Price"])
        ws2.append([1, "Widget", 9.99])
        ws2.append([2, "Gadget", 19.99])

        wb.save(spreadsheets_dir / "sample.xlsx")
        print("Created sample.xlsx")

    except ImportError:
        print("openpyxl not installed, skipping XLSX generation")


def generate_docx_files():
    """Generate DOCX file fixtures using python-docx."""
    documents_dir = FIXTURES_DIR / "documents"
    documents_dir.mkdir(exist_ok=True)

    try:
        from docx import Document
        from docx.shared import Inches

        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is the first paragraph of the test document.")
        doc.add_heading("Section 1", level=1)
        doc.add_paragraph("Content in section 1.")
        doc.add_heading("Section 2", level=1)
        doc.add_paragraph("Content in section 2.")

        # Add a simple table
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Item 1"
        table.cell(1, 1).text = "100"
        table.cell(2, 0).text = "Item 2"
        table.cell(2, 1).text = "200"

        doc.save(documents_dir / "sample.docx")
        print("Created sample.docx")

    except ImportError:
        print("python-docx not installed, skipping DOCX generation")


def generate_pptx_files():
    """Generate PPTX file fixtures using python-pptx."""
    documents_dir = FIXTURES_DIR / "documents"
    documents_dir.mkdir(exist_ok=True)

    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Slide 1 - Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "A test presentation for ingestor"

        # Slide 2 - Content slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        body = slide2.placeholders[1]
        title2.text = "Content Slide"
        tf = body.text_frame
        tf.text = "First bullet point"
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p.level = 1

        prs.save(documents_dir / "sample.pptx")
        print("Created sample.pptx")

    except ImportError:
        print("python-pptx not installed, skipping PPTX generation")


def generate_wav_files():
    """Download public domain speech audio for testing."""
    audio_dir = FIXTURES_DIR / "audio"
    audio_dir.mkdir(exist_ok=True)

    wav_path = audio_dir / "sample.wav"

    # Try multiple sources for public domain speech
    speech_sources = [
        # LibriVox sample (public domain audiobooks)
        "https://ia800301.us.archive.org/5/items/count_monte_cristo_0711_librivox/count_of_monte_cristo_001_dumas_64kb.mp3",
        # Alternative: Short sample from archive.org
        "https://archive.org/download/first_reading_book_librivox/firstreading_01_various.mp3",
    ]

    downloaded = False
    for url in speech_sources:
        try:
            import urllib.request
            print(f"Downloading speech sample from {url[:50]}...")

            # Download to temp file first
            temp_path = audio_dir / "temp_audio"
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=30) as response:
                with open(temp_path, "wb") as f:
                    # Only download first 500KB (enough for a short clip)
                    f.write(response.read(500000))

            # Convert to WAV using ffmpeg if available, otherwise just use the file
            try:
                import subprocess
                # Try to convert with ffmpeg (trim to 5 seconds)
                result = subprocess.run(
                    ["ffmpeg", "-y", "-i", str(temp_path), "-t", "5", "-ar", "16000", "-ac", "1", str(wav_path)],
                    capture_output=True,
                    timeout=30,
                )
                if result.returncode == 0:
                    temp_path.unlink()
                    print("Created sample.wav (public domain speech, 5 seconds)")
                    downloaded = True
                    break
            except (subprocess.SubprocessError, FileNotFoundError):
                # ffmpeg not available, skip this source
                temp_path.unlink()
                continue

        except Exception as e:
            print(f"  Failed: {e}")
            continue

    if not downloaded:
        print("Could not download speech sample. Creating tone (won't test transcription properly)")
        # Fallback to tone if download fails
        import math
        sample_rate = 16000
        duration = 2
        frequency = 440
        samples = []
        for i in range(sample_rate * duration):
            t = i / sample_rate
            value = int(32767 * 0.5 * math.sin(2 * math.pi * frequency * t))
            samples.append(struct.pack("<h", value))
        with wave.open(str(wav_path), "wb") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(b"".join(samples))
        print("Created sample.wav (fallback tone)")


def generate_zip_files():
    """Generate ZIP archive fixtures."""
    import zipfile

    archives_dir = FIXTURES_DIR / "archives"
    archives_dir.mkdir(exist_ok=True)

    zip_path = archives_dir / "sample.zip"
    with zipfile.ZipFile(zip_path, "w") as zf:
        # Add a text file
        zf.writestr("readme.txt", "This is a readme file inside the archive.")
        # Add a JSON file
        zf.writestr(
            "data.json",
            json.dumps({"name": "Archive Test", "value": 42}),
        )
        # Add a nested file
        zf.writestr("folder/nested.txt", "Nested file content.")

    print("Created sample.zip")


def generate_epub_files():
    """Generate EPUB ebook fixtures using ebooklib."""
    ebooks_dir = FIXTURES_DIR / "ebooks"
    ebooks_dir.mkdir(exist_ok=True)

    try:
        from ebooklib import epub

        book = epub.EpubBook()

        # Set metadata
        book.set_identifier("test-book-001")
        book.set_title("Test Ebook")
        book.set_language("en")
        book.add_author("Test Author")

        # Create chapter 1
        c1 = epub.EpubHtml(title="Introduction", file_name="chap_01.xhtml", lang="en")
        c1.content = """
        <html>
        <head><title>Introduction</title></head>
        <body>
            <h1>Introduction</h1>
            <p>This is the introduction chapter of our test ebook.</p>
            <p>It contains multiple paragraphs to test extraction.</p>
        </body>
        </html>
        """

        # Create chapter 2
        c2 = epub.EpubHtml(title="Chapter One", file_name="chap_02.xhtml", lang="en")
        c2.content = """
        <html>
        <head><title>Chapter One</title></head>
        <body>
            <h1>Chapter One</h1>
            <p>This is the first chapter with actual content.</p>
            <h2>Section 1.1</h2>
            <p>Some text in section 1.1.</p>
            <h2>Section 1.2</h2>
            <p>Some text in section 1.2.</p>
        </body>
        </html>
        """

        # Add chapters to book
        book.add_item(c1)
        book.add_item(c2)

        # Define Table of Contents
        book.toc = (
            epub.Link("chap_01.xhtml", "Introduction", "intro"),
            epub.Link("chap_02.xhtml", "Chapter One", "chap1"),
        )

        # Add navigation files
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())

        # Define spine
        book.spine = ["nav", c1, c2]

        # Write EPUB file
        epub_path = ebooks_dir / "sample.epub"
        epub.write_epub(str(epub_path), book)
        print("Created sample.epub")

    except ImportError:
        print("ebooklib not installed, skipping EPUB generation")


def main():
    """Generate all fixture files."""
    print("Generating test fixtures...")
    print("-" * 40)

    generate_text_files()
    generate_json_files()
    generate_xml_files()
    generate_csv_files()
    generate_image_files()
    generate_xlsx_files()
    generate_docx_files()
    generate_pptx_files()
    generate_epub_files()
    generate_wav_files()
    generate_zip_files()

    print("-" * 40)
    print("Done! Fixtures generated in:", FIXTURES_DIR)


if __name__ == "__main__":
    main()
