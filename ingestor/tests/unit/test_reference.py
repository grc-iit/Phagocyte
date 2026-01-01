"""Reference tests for exact output comparison.

These tests compare extractor output against known reference files.
Any difference indicates a potential regression.
"""

import difflib
from pathlib import Path

import pytest

from ingestor.extractors.data.csv_extractor import CsvExtractor
from ingestor.extractors.data.json_extractor import JsonExtractor
from ingestor.extractors.text.txt_extractor import TxtExtractor

# ============================================================================
# Fixtures
# ============================================================================

@pytest.fixture
def reference_dir() -> Path:
    """Path to reference fixtures directory."""
    return Path(__file__).parent.parent / "fixtures" / "reference"


@pytest.fixture
def txt_extractor():
    """Text extractor instance."""
    return TxtExtractor()


@pytest.fixture
def json_extractor():
    """JSON extractor instance."""
    return JsonExtractor()


@pytest.fixture
def csv_extractor():
    """CSV extractor instance."""
    return CsvExtractor()


# ============================================================================
# Helper Functions
# ============================================================================

def normalize_markdown(text: str) -> str:
    """Normalize markdown for comparison.

    - Strip trailing whitespace from lines
    - Normalize line endings
    - Strip leading/trailing whitespace from whole text
    """
    lines = text.replace('\r\n', '\n').split('\n')
    lines = [line.rstrip() for line in lines]
    return '\n'.join(lines).strip()


def compare_markdown(actual: str, expected: str) -> tuple[bool, str]:
    """Compare two markdown strings and return diff if different.

    Returns:
        (match, diff_output) - True if match, otherwise diff string
    """
    actual_norm = normalize_markdown(actual)
    expected_norm = normalize_markdown(expected)

    if actual_norm == expected_norm:
        return True, ""

    # Generate unified diff
    actual_lines = actual_norm.split('\n')
    expected_lines = expected_norm.split('\n')

    diff = difflib.unified_diff(
        expected_lines,
        actual_lines,
        fromfile='expected',
        tofile='actual',
        lineterm=''
    )

    return False, '\n'.join(diff)


def load_reference(reference_dir: Path, input_name: str) -> tuple[Path | None, str | None]:
    """Load a reference input file and its expected output.

    Args:
        reference_dir: Directory containing reference files
        input_name: Name of input file (e.g., "simple.md")

    Returns:
        (input_path, expected_content) or (None, None) if not found
    """
    input_path = reference_dir / input_name
    expected_path = reference_dir / f"{input_name}.expected.md"

    if not input_path.exists():
        return None, None

    if not expected_path.exists():
        return input_path, None

    expected_content = expected_path.read_text(encoding='utf-8')
    return input_path, expected_content


# ============================================================================
# Text Reference Tests
# ============================================================================

class TestTextReferences:
    """Reference tests for text file extraction."""

    @pytest.mark.asyncio
    async def test_simple_md_reference(self, reference_dir, txt_extractor):
        """Test simple markdown extraction matches reference."""
        input_path, expected = load_reference(reference_dir, "simple.md")

        if input_path is None:
            pytest.skip("Reference file simple.md not found")

        result = await txt_extractor.extract(input_path)

        if expected is not None:
            match, diff = compare_markdown(result.markdown, expected)
            assert match, f"Output doesn't match reference:\n{diff}"

    @pytest.mark.asyncio
    async def test_empty_file_reference(self, reference_dir, txt_extractor):
        """Test empty file extraction."""
        input_path, expected = load_reference(reference_dir, "empty.txt")

        if input_path is None:
            pytest.skip("Reference file empty.txt not found")

        result = await txt_extractor.extract(input_path)

        # Empty file should produce empty or minimal output
        assert len(result.markdown.strip()) == 0 or result.markdown.strip() == ""

    @pytest.mark.asyncio
    async def test_whitespace_only_reference(self, reference_dir, txt_extractor):
        """Test whitespace-only file extraction."""
        input_path, expected = load_reference(reference_dir, "whitespace_only.txt")

        if input_path is None:
            pytest.skip("Reference file whitespace_only.txt not found")

        result = await txt_extractor.extract(input_path)

        # Whitespace-only should produce empty output
        assert len(result.markdown.strip()) == 0

    @pytest.mark.asyncio
    async def test_unicode_content_reference(self, reference_dir, txt_extractor):
        """Test unicode content extraction."""
        input_path, expected = load_reference(reference_dir, "unicode_content.txt")

        if input_path is None:
            pytest.skip("Reference file unicode_content.txt not found")

        result = await txt_extractor.extract(input_path)

        # Check unicode characters are preserved
        assert "こんにちは" in result.markdown
        assert "مرحبا" in result.markdown
        assert "Привет" in result.markdown
        assert "🎉" in result.markdown

    @pytest.mark.asyncio
    async def test_special_chars_reference(self, reference_dir, txt_extractor):
        """Test special characters extraction."""
        input_path, expected = load_reference(reference_dir, "special_chars.txt")

        if input_path is None:
            pytest.skip("Reference file special_chars.txt not found")

        result = await txt_extractor.extract(input_path)

        if expected is not None:
            match, diff = compare_markdown(result.markdown, expected)
            assert match, f"Output doesn't match reference:\n{diff}"

    @pytest.mark.asyncio
    async def test_long_line_reference(self, reference_dir, txt_extractor):
        """Test long line handling."""
        input_path, expected = load_reference(reference_dir, "long_line.txt")

        if input_path is None:
            pytest.skip("Reference file long_line.txt not found")

        result = await txt_extractor.extract(input_path)

        # Should preserve long lines
        assert "A" * 100 in result.markdown  # At least 100 A's


# ============================================================================
# JSON Reference Tests
# ============================================================================

class TestJsonReferences:
    """Reference tests for JSON file extraction."""

    @pytest.mark.asyncio
    async def test_edge_cases_json(self, reference_dir, json_extractor):
        """Test JSON edge cases extraction."""
        input_path, expected = load_reference(reference_dir, "edge_cases.json")

        if input_path is None:
            pytest.skip("Reference file edge_cases.json not found")

        result = await json_extractor.extract(input_path)

        # Check key content is preserved
        assert "null" in result.markdown.lower() or "None" in result.markdown
        assert "日本語テスト" in result.markdown
        assert "42" in result.markdown
        assert "3.14159" in result.markdown


# ============================================================================
# CSV Reference Tests
# ============================================================================

class TestCsvReferences:
    """Reference tests for CSV file extraction."""

    @pytest.mark.asyncio
    async def test_edge_cases_csv(self, reference_dir, csv_extractor):
        """Test CSV edge cases extraction."""
        input_path, expected = load_reference(reference_dir, "edge_cases.csv")

        if input_path is None:
            pytest.skip("Reference file edge_cases.csv not found")

        result = await csv_extractor.extract(input_path)

        # Check table structure exists
        assert "|" in result.markdown  # Table format
        assert "Simple" in result.markdown
        assert "With, Comma" in result.markdown or "With Comma" in result.markdown


# ============================================================================
# Document Reference Tests (require optional dependencies)
# ============================================================================

class TestDocxReferences:
    """Reference tests for DOCX extraction."""

    @pytest.fixture
    def docx_extractor(self):
        """DOCX extractor instance."""
        try:
            from ingestor.extractors.docx.docx_extractor import DocxExtractor
            return DocxExtractor()
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_complex_docx_reference(self, reference_dir, docx_extractor):
        """Test complex DOCX extraction."""
        input_path, expected = load_reference(reference_dir, "complex.docx")

        if input_path is None:
            pytest.skip("Reference file complex.docx not found")

        result = await docx_extractor.extract(input_path)

        # Check key content is extracted
        assert "Complex Document Test" in result.markdown
        assert "Section 1" in result.markdown
        assert "Tables" in result.markdown
        assert "Alice" in result.markdown
        assert "Lists" in result.markdown


class TestXlsxReferences:
    """Reference tests for XLSX extraction."""

    @pytest.fixture
    def xlsx_extractor(self):
        """XLSX extractor instance."""
        try:
            from ingestor.extractors.excel.xlsx_extractor import XlsxExtractor
            return XlsxExtractor()
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_complex_xlsx_reference(self, reference_dir, xlsx_extractor):
        """Test complex XLSX extraction."""
        input_path, expected = load_reference(reference_dir, "complex.xlsx")

        if input_path is None:
            pytest.skip("Reference file complex.xlsx not found")

        result = await xlsx_extractor.extract(input_path)

        # Check key content is extracted
        assert "Sales" in result.markdown or "Product" in result.markdown
        assert "Widget" in result.markdown
        assert "|" in result.markdown  # Table format


class TestPptxReferences:
    """Reference tests for PPTX extraction."""

    @pytest.fixture
    def pptx_extractor(self):
        """PPTX extractor instance."""
        try:
            from ingestor.extractors.pptx.pptx_extractor import PptxExtractor
            return PptxExtractor()
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_complex_pptx_reference(self, reference_dir, pptx_extractor):
        """Test complex PPTX extraction."""
        input_path, expected = load_reference(reference_dir, "complex.pptx")

        if input_path is None:
            pytest.skip("Reference file complex.pptx not found")

        result = await pptx_extractor.extract(input_path)

        # Check key content is extracted
        assert "Complex Presentation Test" in result.markdown
        assert "Bullet" in result.markdown
        assert "Table" in result.markdown or "Data" in result.markdown


# ============================================================================
# Additional Complex Feature Tests
# ============================================================================

class TestDocxAdvancedFeatures:
    """Tests for advanced DOCX features."""

    @pytest.fixture
    def docx_extractor(self):
        try:
            from ingestor.extractors.docx.docx_extractor import DocxExtractor
            return DocxExtractor()
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_docx_multiple_sections(self, docx_extractor, tmp_path):
        """Test DOCX with multiple sections and page breaks."""
        try:
            from docx import Document
            from docx.enum.section import WD_ORIENT

            doc = Document()

            # Section 1
            doc.add_heading("Section 1", level=1)
            doc.add_paragraph("Content in first section.")

            # Add section break
            doc.add_section()

            # Section 2
            doc.add_heading("Section 2", level=1)
            doc.add_paragraph("Content in second section.")

            doc.save(tmp_path / "multi_section.docx")

            result = await docx_extractor.extract(tmp_path / "multi_section.docx")

            assert "Section 1" in result.markdown
            assert "Section 2" in result.markdown
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_docx_footnotes_endnotes(self, docx_extractor, tmp_path):
        """Test DOCX content even without explicit footnote support."""
        try:
            from docx import Document

            doc = Document()
            doc.add_heading("Document with References", level=1)
            doc.add_paragraph("Main text with inline references [1].")
            doc.add_paragraph("More content with citations [2].")
            doc.add_heading("References", level=2)
            doc.add_paragraph("[1] First reference source")
            doc.add_paragraph("[2] Second reference source")

            doc.save(tmp_path / "with_refs.docx")

            result = await docx_extractor.extract(tmp_path / "with_refs.docx")

            assert "References" in result.markdown
            assert "[1]" in result.markdown
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_docx_track_changes_simulation(self, docx_extractor, tmp_path):
        """Test that documents simulate track changes are processed."""
        try:
            from docx import Document
            from docx.shared import RGBColor

            doc = Document()
            doc.add_heading("Document with Revisions", level=1)

            # Simulate tracked change with strikethrough/color
            para = doc.add_paragraph()
            para.add_run("Original text ")
            run2 = para.add_run("deleted text")
            run2.font.strike = True
            run2.font.color.rgb = RGBColor(255, 0, 0)
            run3 = para.add_run(" new text")
            run3.font.color.rgb = RGBColor(0, 128, 0)

            doc.add_paragraph("Final paragraph without changes.")

            doc.save(tmp_path / "track_changes.docx")

            result = await docx_extractor.extract(tmp_path / "track_changes.docx")

            assert result is not None
            assert "Document with Revisions" in result.markdown
        except ImportError:
            pytest.skip("python-docx not installed")


class TestXlsxAdvancedFeatures:
    """Tests for advanced XLSX features."""

    @pytest.fixture
    def xlsx_extractor(self):
        try:
            from ingestor.extractors.excel.xlsx_extractor import XlsxExtractor
            return XlsxExtractor()
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_xlsx_conditional_formatting(self, xlsx_extractor, tmp_path):
        """Test XLSX with conditional formatting extracts data."""
        try:
            from openpyxl import Workbook
            from openpyxl.formatting.rule import ColorScaleRule

            wb = Workbook()
            ws = wb.active

            ws["A1"] = "Score"
            for i in range(2, 12):
                ws[f"A{i}"] = i * 10

            # Add conditional formatting
            ws.conditional_formatting.add(
                "A2:A11",
                ColorScaleRule(
                    start_type="min", start_color="FF0000",
                    end_type="max", end_color="00FF00"
                )
            )

            wb.save(tmp_path / "conditional.xlsx")

            result = await xlsx_extractor.extract(tmp_path / "conditional.xlsx")

            assert result is not None
            assert "Score" in result.markdown
            assert "100" in result.markdown
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_xlsx_data_validation(self, xlsx_extractor, tmp_path):
        """Test XLSX with data validation extracts data."""
        try:
            from openpyxl import Workbook
            from openpyxl.worksheet.datavalidation import DataValidation

            wb = Workbook()
            ws = wb.active

            ws["A1"] = "Status"
            ws["A2"] = "Active"
            ws["A3"] = "Inactive"

            # Add data validation (dropdown)
            dv = DataValidation(
                type="list",
                formula1='"Active,Inactive,Pending"',
                allow_blank=True
            )
            ws.add_data_validation(dv)
            dv.add(ws["A2"])
            dv.add(ws["A3"])

            wb.save(tmp_path / "validation.xlsx")

            result = await xlsx_extractor.extract(tmp_path / "validation.xlsx")

            assert result is not None
            assert "Status" in result.markdown
            assert "Active" in result.markdown
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_xlsx_hidden_sheets(self, xlsx_extractor, tmp_path):
        """Test XLSX with hidden sheets."""
        try:
            from openpyxl import Workbook

            wb = Workbook()

            # Visible sheet with proper table data
            ws1 = wb.active
            ws1.title = "Visible"
            ws1["A1"] = "Name"
            ws1["B1"] = "Value"
            ws1["A2"] = "Visible Data"
            ws1["B2"] = 100

            # Hidden sheet
            ws2 = wb.create_sheet("Hidden")
            ws2["A1"] = "Hidden Data"
            ws2.sheet_state = 'hidden'

            wb.save(tmp_path / "hidden_sheet.xlsx")

            result = await xlsx_extractor.extract(tmp_path / "hidden_sheet.xlsx")

            assert result is not None
            # Check that visible sheet data is extracted
            assert "Visible" in result.markdown or "Name" in result.markdown or "100" in result.markdown
        except ImportError:
            pytest.skip("openpyxl not installed")


class TestPptxAdvancedFeatures:
    """Tests for advanced PPTX features."""

    @pytest.fixture
    def pptx_extractor(self):
        try:
            from ingestor.extractors.pptx.pptx_extractor import PptxExtractor
            return PptxExtractor()
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_pptx_slide_transitions(self, pptx_extractor, tmp_path):
        """Test PPTX with slide transitions extracts content."""
        try:
            from pptx import Presentation

            prs = Presentation()

            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "Slide with Transition"
            slide1.shapes.placeholders[1].text = "Content on slide 1"

            slide2 = prs.slides.add_slide(prs.slide_layouts[0])
            slide2.shapes.title.text = "Another Slide"
            slide2.shapes.placeholders[1].text = "Content on slide 2"

            prs.save(tmp_path / "transitions.pptx")

            result = await pptx_extractor.extract(tmp_path / "transitions.pptx")

            assert result is not None
            assert "Slide with Transition" in result.markdown
            assert "Another Slide" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_pptx_master_slide_content(self, pptx_extractor, tmp_path):
        """Test PPTX extracts slide content regardless of master."""
        try:
            from pptx import Presentation

            prs = Presentation()

            # Use different layouts
            for _i, layout_idx in enumerate([0, 1, 5, 6]):
                try:
                    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                    if slide.shapes.title:
                        slide.shapes.title.text = f"Layout {layout_idx} Title"
                except:
                    pass

            prs.save(tmp_path / "layouts.pptx")

            result = await pptx_extractor.extract(tmp_path / "layouts.pptx")

            assert result is not None
            assert "Layout" in result.markdown or "Title" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_pptx_grouped_shapes(self, pptx_extractor, tmp_path):
        """Test PPTX with grouped shapes extracts text."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add individual shapes (grouping requires more complex API)
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(1), Inches(2), Inches(1)
            )
            shape1.text = "Shape 1 Text"

            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(2.5), Inches(2), Inches(1)
            )
            shape2.text = "Shape 2 Text"

            prs.save(tmp_path / "shapes.pptx")

            result = await pptx_extractor.extract(tmp_path / "shapes.pptx")

            assert result is not None
            assert "Shape 1 Text" in result.markdown
            assert "Shape 2 Text" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")
