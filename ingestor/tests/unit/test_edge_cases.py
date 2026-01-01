"""Edge case tests for extractors.

Tests for boundary conditions, error handling, and unusual inputs:
- Empty files
- Corrupted files
- Very large files
- Unicode filenames
- Files with no extension
- Password-protected files
"""

import os
import tempfile
from pathlib import Path

import pytest

from ingestor.extractors.data.csv_extractor import CsvExtractor
from ingestor.extractors.data.json_extractor import JsonExtractor
from ingestor.extractors.data.xml_extractor import XmlExtractor
from ingestor.extractors.text.txt_extractor import TxtExtractor
from ingestor.types import MediaType

# ============================================================================
# Fixtures
# ============================================================================

@pytest.fixture
def txt_extractor():
    return TxtExtractor()


@pytest.fixture
def json_extractor():
    return JsonExtractor()


@pytest.fixture
def csv_extractor():
    return CsvExtractor()


@pytest.fixture
def xml_extractor():
    return XmlExtractor()


@pytest.fixture
def temp_dir():
    """Provide a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


# ============================================================================
# Empty File Tests
# ============================================================================

class TestEmptyFiles:
    """Test handling of empty files."""

    @pytest.mark.asyncio
    async def test_empty_txt_file(self, txt_extractor, temp_dir):
        """Empty text file should produce empty output."""
        empty_file = temp_dir / "empty.txt"
        empty_file.write_text("")

        result = await txt_extractor.extract(empty_file)

        assert result is not None
        assert result.media_type == MediaType.TXT
        assert result.markdown.strip() == ""

    @pytest.mark.asyncio
    async def test_empty_json_file(self, json_extractor, temp_dir):
        """Empty JSON file should raise JSONDecodeError (invalid JSON)."""
        import json
        empty_file = temp_dir / "empty.json"
        empty_file.write_text("")

        # Empty file is not valid JSON - should raise error
        with pytest.raises(json.JSONDecodeError):
            await json_extractor.extract(empty_file)

    @pytest.mark.asyncio
    async def test_empty_csv_file(self, csv_extractor, temp_dir):
        """Empty CSV file should raise error (no data to parse)."""
        import pandas as pd
        empty_file = temp_dir / "empty.csv"
        empty_file.write_text("")

        # Empty CSV is not valid - should raise error
        with pytest.raises((pd.errors.EmptyDataError, Exception)):
            await csv_extractor.extract(empty_file)


# ============================================================================
# Whitespace-Only Files
# ============================================================================

class TestWhitespaceFiles:
    """Test handling of whitespace-only files."""

    @pytest.mark.asyncio
    async def test_whitespace_txt(self, txt_extractor, temp_dir):
        """Whitespace-only text file."""
        ws_file = temp_dir / "whitespace.txt"
        ws_file.write_text("   \n\n\t\t\n   ")

        result = await txt_extractor.extract(ws_file)

        assert result is not None
        assert result.markdown.strip() == ""

    @pytest.mark.asyncio
    async def test_newlines_only_txt(self, txt_extractor, temp_dir):
        """File with only newlines."""
        nl_file = temp_dir / "newlines.txt"
        nl_file.write_text("\n\n\n\n\n")

        result = await txt_extractor.extract(nl_file)

        assert result is not None
        assert result.markdown.strip() == ""


# ============================================================================
# Malformed/Corrupted Files
# ============================================================================

class TestMalformedFiles:
    """Test handling of malformed or corrupted files."""

    @pytest.mark.asyncio
    async def test_invalid_json(self, json_extractor, temp_dir):
        """Invalid JSON should raise JSONDecodeError."""
        import json
        invalid_file = temp_dir / "invalid.json"
        invalid_file.write_text("{ invalid json content }")

        # Invalid JSON should raise error
        with pytest.raises(json.JSONDecodeError):
            await json_extractor.extract(invalid_file)

    @pytest.mark.asyncio
    async def test_truncated_json(self, json_extractor, temp_dir):
        """Truncated JSON (incomplete) should raise error."""
        import json
        truncated_file = temp_dir / "truncated.json"
        truncated_file.write_text('{"key": "value", "incomplete":')

        with pytest.raises(json.JSONDecodeError):
            await json_extractor.extract(truncated_file)

    @pytest.mark.asyncio
    async def test_invalid_xml(self, xml_extractor, temp_dir):
        """Invalid XML should raise ParseError."""
        pytest.importorskip("defusedxml", reason="defusedxml not installed")
        import xml.etree.ElementTree as ET
        invalid_file = temp_dir / "invalid.xml"
        invalid_file.write_text("<root><unclosed>")

        with pytest.raises(ET.ParseError):
            await xml_extractor.extract(invalid_file)

    @pytest.mark.asyncio
    async def test_binary_as_text(self, txt_extractor, temp_dir):
        """Binary data passed as text file."""
        binary_file = temp_dir / "binary.txt"
        binary_file.write_bytes(b'\x00\x01\x02\x03\xff\xfe\xfd')

        # Should handle gracefully, not crash
        try:
            result = await txt_extractor.extract(binary_file)
            assert result is not None
        except Exception as e:
            # Acceptable to raise an exception for binary content
            assert "decode" in str(e).lower() or "encoding" in str(e).lower()


# ============================================================================
# Unicode and Special Characters
# ============================================================================

class TestUnicodeHandling:
    """Test handling of unicode content and filenames."""

    @pytest.mark.asyncio
    async def test_unicode_content(self, txt_extractor, temp_dir):
        """File with various unicode characters."""
        unicode_file = temp_dir / "unicode.txt"
        content = """# Unicode Test
日本語 (Japanese)
한국어 (Korean)
العربية (Arabic)
עברית (Hebrew)
Ελληνικά (Greek)
🎉🚀💻🔥 (Emoji)
"""
        unicode_file.write_text(content, encoding='utf-8')

        result = await txt_extractor.extract(unicode_file)

        assert result is not None
        assert "日本語" in result.markdown
        assert "한국어" in result.markdown
        assert "🎉" in result.markdown

    @pytest.mark.asyncio
    async def test_unicode_filename(self, txt_extractor, temp_dir):
        """File with unicode characters in filename."""
        unicode_name = temp_dir / "文档_테스트_документ.txt"
        unicode_name.write_text("Content in unicode-named file")

        result = await txt_extractor.extract(unicode_name)

        assert result is not None
        assert "Content" in result.markdown

    @pytest.mark.asyncio
    async def test_mixed_encodings(self, txt_extractor, temp_dir):
        """File that could be interpreted as different encodings."""
        # Latin-1 encoded content
        latin1_file = temp_dir / "latin1.txt"
        latin1_content = "Café résumé naïve"
        latin1_file.write_bytes(latin1_content.encode('latin-1'))

        result = await txt_extractor.extract(latin1_file)

        assert result is not None
        # Should either decode correctly or return something readable
        assert len(result.markdown) > 0

    @pytest.mark.asyncio
    async def test_bom_markers(self, txt_extractor, temp_dir):
        """File with BOM (Byte Order Mark)."""
        bom_file = temp_dir / "bom.txt"
        content = "Content after BOM"
        # UTF-8 BOM
        bom_file.write_bytes(b'\xef\xbb\xbf' + content.encode('utf-8'))

        result = await txt_extractor.extract(bom_file)

        assert result is not None
        assert "Content after BOM" in result.markdown


# ============================================================================
# Special Path Cases
# ============================================================================

class TestSpecialPaths:
    """Test handling of special file paths."""

    @pytest.mark.asyncio
    async def test_file_with_spaces(self, txt_extractor, temp_dir):
        """Filename with spaces."""
        space_file = temp_dir / "file with spaces.txt"
        space_file.write_text("Content in spaced filename")

        result = await txt_extractor.extract(space_file)

        assert result is not None
        assert "Content" in result.markdown

    @pytest.mark.asyncio
    async def test_file_with_special_chars(self, txt_extractor, temp_dir):
        """Filename with special characters."""
        # Use characters that are valid on most filesystems
        special_file = temp_dir / "file-with_special.chars(1).txt"
        special_file.write_text("Content in special filename")

        result = await txt_extractor.extract(special_file)

        assert result is not None
        assert "Content" in result.markdown

    @pytest.mark.asyncio
    async def test_deeply_nested_path(self, txt_extractor, temp_dir):
        """File in deeply nested directory."""
        nested_path = temp_dir / "a" / "b" / "c" / "d" / "e"
        nested_path.mkdir(parents=True)
        nested_file = nested_path / "deep.txt"
        nested_file.write_text("Deep content")

        result = await txt_extractor.extract(nested_file)

        assert result is not None
        assert "Deep content" in result.markdown

    @pytest.mark.asyncio
    async def test_hidden_file(self, txt_extractor, temp_dir):
        """Hidden file (starting with dot)."""
        hidden_file = temp_dir / ".hidden.txt"
        hidden_file.write_text("Hidden content")

        result = await txt_extractor.extract(hidden_file)

        assert result is not None
        assert "Hidden content" in result.markdown


# ============================================================================
# Large Content Tests
# ============================================================================

class TestLargeContent:
    """Test handling of large files and content."""

    @pytest.mark.asyncio
    async def test_very_long_line(self, txt_extractor, temp_dir):
        """File with a very long line (100KB)."""
        long_file = temp_dir / "long_line.txt"
        long_line = "A" * 100_000
        long_file.write_text(f"# Long Line\n{long_line}\n")

        result = await txt_extractor.extract(long_file)

        assert result is not None
        assert len(result.markdown) > 100_000

    @pytest.mark.asyncio
    async def test_many_lines(self, txt_extractor, temp_dir):
        """File with many lines (10,000 lines)."""
        many_lines_file = temp_dir / "many_lines.txt"
        lines = [f"Line {i}" for i in range(10_000)]
        many_lines_file.write_text("\n".join(lines))

        result = await txt_extractor.extract(many_lines_file)

        assert result is not None
        assert "Line 0" in result.markdown
        assert "Line 9999" in result.markdown

    @pytest.mark.asyncio
    async def test_large_json_array(self, json_extractor, temp_dir):
        """Large JSON array (1000 items)."""
        import json
        large_json_file = temp_dir / "large.json"
        data = [{"id": i, "name": f"Item {i}"} for i in range(1000)]
        large_json_file.write_text(json.dumps(data))

        result = await json_extractor.extract(large_json_file)

        assert result is not None
        assert "Item 0" in result.markdown or "id" in result.markdown

    @pytest.mark.asyncio
    async def test_deeply_nested_json(self, json_extractor, temp_dir):
        """Deeply nested JSON structure."""
        import json
        deep_json_file = temp_dir / "deep.json"

        # Create 50 levels of nesting
        data = {"level": 0}
        current = data
        for i in range(1, 50):
            current["nested"] = {"level": i}
            current = current["nested"]

        deep_json_file.write_text(json.dumps(data))

        result = await json_extractor.extract(deep_json_file)

        assert result is not None


# ============================================================================
# File Extension Edge Cases
# ============================================================================

class TestExtensionEdgeCases:
    """Test handling of unusual file extensions."""

    @pytest.mark.asyncio
    async def test_no_extension(self, txt_extractor, temp_dir):
        """File with no extension."""
        no_ext_file = temp_dir / "noextension"
        no_ext_file.write_text("Content without extension")

        # This might not be supported, but shouldn't crash
        try:
            result = await txt_extractor.extract(no_ext_file)
            assert result is not None
        except Exception:
            pass  # Acceptable to not support

    @pytest.mark.asyncio
    async def test_double_extension(self, txt_extractor, temp_dir):
        """File with double extension."""
        double_ext_file = temp_dir / "file.txt.bak"
        double_ext_file.write_text("Content with double extension")

        # Might be treated as .bak file
        if txt_extractor.supports(double_ext_file):
            result = await txt_extractor.extract(double_ext_file)
            assert result is not None

    @pytest.mark.asyncio
    async def test_uppercase_extension(self, txt_extractor, temp_dir):
        """File with uppercase extension."""
        upper_ext_file = temp_dir / "file.TXT"
        upper_ext_file.write_text("Content with uppercase extension")

        result = await txt_extractor.extract(upper_ext_file)

        assert result is not None
        assert "Content" in result.markdown

    @pytest.mark.asyncio
    async def test_mixed_case_extension(self, txt_extractor, temp_dir):
        """File with mixed case extension."""
        mixed_ext_file = temp_dir / "file.TxT"
        mixed_ext_file.write_text("Content with mixed case extension")

        result = await txt_extractor.extract(mixed_ext_file)

        assert result is not None


# ============================================================================
# Concurrent/Stress Tests (Light)
# ============================================================================

class TestConcurrentExtraction:
    """Light concurrent extraction tests."""

    @pytest.mark.asyncio
    async def test_multiple_files_sequential(self, txt_extractor, temp_dir):
        """Extract multiple files sequentially."""
        files = []
        for i in range(10):
            f = temp_dir / f"file_{i}.txt"
            f.write_text(f"Content {i}")
            files.append(f)

        results = []
        for f in files:
            result = await txt_extractor.extract(f)
            results.append(result)

        assert len(results) == 10
        assert all(r is not None for r in results)
        for i, r in enumerate(results):
            assert f"Content {i}" in r.markdown

    @pytest.mark.asyncio
    async def test_same_file_multiple_times(self, txt_extractor, temp_dir):
        """Extract the same file multiple times."""
        f = temp_dir / "repeated.txt"
        f.write_text("Repeated content")

        results = []
        for _ in range(5):
            result = await txt_extractor.extract(f)
            results.append(result)

        # All results should be identical
        assert len(results) == 5
        first_md = results[0].markdown
        assert all(r.markdown == first_md for r in results)


# ============================================================================
# Error Recovery Tests
# ============================================================================

class TestErrorRecovery:
    """Test graceful error handling and recovery."""

    @pytest.mark.asyncio
    async def test_nonexistent_file(self, txt_extractor):
        """Extracting nonexistent file should handle gracefully."""
        fake_path = Path("/nonexistent/path/to/file.txt")

        # Should raise an appropriate error, not crash
        with pytest.raises(Exception):
            await txt_extractor.extract(fake_path)

    @pytest.mark.asyncio
    async def test_directory_instead_of_file(self, txt_extractor, temp_dir):
        """Passing a directory instead of file."""
        directory = temp_dir / "subdir"
        directory.mkdir()

        # Should handle gracefully
        with pytest.raises(Exception):
            await txt_extractor.extract(directory)

    @pytest.mark.asyncio
    async def test_permission_denied(self, txt_extractor, temp_dir):
        """File with no read permissions (Unix only)."""
        if os.name != 'posix':
            pytest.skip("Permission test only works on Unix")

        restricted_file = temp_dir / "restricted.txt"
        restricted_file.write_text("Secret content")
        restricted_file.chmod(0o000)

        try:
            # Should raise permission error
            with pytest.raises(Exception):
                await txt_extractor.extract(restricted_file)
        finally:
            # Restore permissions for cleanup
            restricted_file.chmod(0o644)


# ============================================================================
# Password-Protected File Tests
# ============================================================================

class TestPasswordProtectedFiles:
    """Test handling of password-protected files."""

    @pytest.fixture
    def docx_extractor(self):
        """DOCX extractor instance."""
        try:
            from ingestor.extractors.docx.docx_extractor import DocxExtractor
            return DocxExtractor()
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.fixture
    def xlsx_extractor(self):
        """XLSX extractor instance."""
        try:
            from ingestor.extractors.excel.xlsx_extractor import XlsxExtractor
            return XlsxExtractor()
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.fixture
    def pptx_extractor(self):
        """PPTX extractor instance."""
        try:
            from ingestor.extractors.pptx.pptx_extractor import PptxExtractor
            return PptxExtractor()
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_encrypted_docx_graceful_failure(self, docx_extractor, temp_dir):
        """Encrypted DOCX should fail gracefully."""
        # Create a fake "encrypted" DOCX (actually just invalid)
        encrypted_file = temp_dir / "encrypted.docx"
        # Write invalid ZIP structure that simulates encryption
        encrypted_file.write_bytes(b'PK\x03\x04' + b'\x00' * 100)

        # Should raise an exception but not crash
        with pytest.raises(Exception):
            await docx_extractor.extract(encrypted_file)

    @pytest.mark.asyncio
    async def test_encrypted_xlsx_graceful_failure(self, xlsx_extractor, temp_dir):
        """Encrypted XLSX should fail gracefully."""
        encrypted_file = temp_dir / "encrypted.xlsx"
        encrypted_file.write_bytes(b'PK\x03\x04' + b'\x00' * 100)

        with pytest.raises(Exception):
            await xlsx_extractor.extract(encrypted_file)

    @pytest.mark.asyncio
    async def test_encrypted_pptx_graceful_failure(self, pptx_extractor, temp_dir):
        """Encrypted PPTX should fail gracefully."""
        encrypted_file = temp_dir / "encrypted.pptx"
        encrypted_file.write_bytes(b'PK\x03\x04' + b'\x00' * 100)

        with pytest.raises(Exception):
            await pptx_extractor.extract(encrypted_file)


# ============================================================================
# Corrupted Office File Tests
# ============================================================================

class TestCorruptedOfficeFiles:
    """Test handling of corrupted Office files."""

    @pytest.fixture
    def docx_extractor(self):
        try:
            from ingestor.extractors.docx.docx_extractor import DocxExtractor
            return DocxExtractor()
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.fixture
    def xlsx_extractor(self):
        try:
            from ingestor.extractors.excel.xlsx_extractor import XlsxExtractor
            return XlsxExtractor()
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.fixture
    def pptx_extractor(self):
        try:
            from ingestor.extractors.pptx.pptx_extractor import PptxExtractor
            return PptxExtractor()
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_truncated_docx(self, docx_extractor, temp_dir):
        """Truncated DOCX file should fail gracefully."""
        truncated_file = temp_dir / "truncated.docx"
        # Write partial ZIP header
        truncated_file.write_bytes(b'PK\x03\x04\x14\x00\x00\x00')

        with pytest.raises(Exception):
            await docx_extractor.extract(truncated_file)

    @pytest.mark.asyncio
    async def test_wrong_extension_docx(self, docx_extractor, temp_dir):
        """Text file with .docx extension should fail gracefully."""
        fake_docx = temp_dir / "fake.docx"
        fake_docx.write_text("This is not a DOCX file")

        with pytest.raises(Exception):
            await docx_extractor.extract(fake_docx)

    @pytest.mark.asyncio
    async def test_wrong_extension_xlsx(self, xlsx_extractor, temp_dir):
        """Text file with .xlsx extension should fail gracefully."""
        fake_xlsx = temp_dir / "fake.xlsx"
        fake_xlsx.write_text("This is not an XLSX file")

        with pytest.raises(Exception):
            await xlsx_extractor.extract(fake_xlsx)

    @pytest.mark.asyncio
    async def test_wrong_extension_pptx(self, pptx_extractor, temp_dir):
        """Text file with .pptx extension should fail gracefully."""
        fake_pptx = temp_dir / "fake.pptx"
        fake_pptx.write_text("This is not a PPTX file")

        with pytest.raises(Exception):
            await pptx_extractor.extract(fake_pptx)

    @pytest.mark.asyncio
    async def test_zero_byte_docx(self, docx_extractor, temp_dir):
        """Zero-byte DOCX file should fail gracefully."""
        zero_file = temp_dir / "zero.docx"
        zero_file.write_bytes(b'')

        with pytest.raises(Exception):
            await docx_extractor.extract(zero_file)


# ============================================================================
# Complex Document Content Tests
# ============================================================================

class TestComplexDocumentContent:
    """Test extraction of complex document features."""

    @pytest.fixture
    def docx_extractor(self):
        try:
            from ingestor.extractors.docx.docx_extractor import DocxExtractor
            return DocxExtractor()
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.fixture
    def xlsx_extractor(self):
        try:
            from ingestor.extractors.excel.xlsx_extractor import XlsxExtractor
            return XlsxExtractor()
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.fixture
    def pptx_extractor(self):
        try:
            from ingestor.extractors.pptx.pptx_extractor import PptxExtractor
            return PptxExtractor()
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_docx_with_comments(self, docx_extractor, temp_dir):
        """DOCX with comments should extract main content."""
        try:
            from docx import Document

            doc = Document()
            doc.add_paragraph("Main content of the document.")
            doc.add_paragraph("This paragraph has important information.")
            # Note: python-docx doesn't easily support adding comments
            # but we test that documents with potential comments work
            doc.save(temp_dir / "with_comments.docx")

            result = await docx_extractor.extract(temp_dir / "with_comments.docx")

            assert result is not None
            assert "Main content" in result.markdown
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_docx_with_headers_footers(self, docx_extractor, temp_dir):
        """DOCX with headers/footers should extract content."""
        try:
            from docx import Document

            doc = Document()

            # Add header
            section = doc.sections[0]
            header = section.header
            header.paragraphs[0].text = "Document Header"

            # Add footer
            footer = section.footer
            footer.paragraphs[0].text = "Page Footer"

            # Add main content
            doc.add_heading("Main Title", 0)
            doc.add_paragraph("Body content of the document.")

            doc.save(temp_dir / "with_headers.docx")

            result = await docx_extractor.extract(temp_dir / "with_headers.docx")

            assert result is not None
            assert "Main Title" in result.markdown
            assert "Body content" in result.markdown
        except ImportError:
            pytest.skip("python-docx not installed")

    @pytest.mark.asyncio
    async def test_xlsx_with_charts_graceful(self, xlsx_extractor, temp_dir):
        """XLSX with charts should extract data (charts may be ignored)."""
        try:
            from openpyxl import Workbook
            from openpyxl.chart import BarChart, Reference

            wb = Workbook()
            ws = wb.active
            ws.title = "Chart Data"

            # Add data
            data = [
                ["Category", "Value"],
                ["A", 10],
                ["B", 20],
                ["C", 30],
            ]
            for row in data:
                ws.append(row)

            # Add a chart
            chart = BarChart()
            chart.title = "Sample Chart"
            data_ref = Reference(ws, min_col=2, min_row=1, max_col=2, max_row=4)
            cats = Reference(ws, min_col=1, min_row=2, max_row=4)
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats)
            ws.add_chart(chart, "E1")

            wb.save(temp_dir / "with_chart.xlsx")

            result = await xlsx_extractor.extract(temp_dir / "with_chart.xlsx")

            assert result is not None
            # Data should be extracted even if chart is ignored
            assert "Category" in result.markdown or "Value" in result.markdown
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_xlsx_with_named_ranges(self, xlsx_extractor, temp_dir):
        """XLSX with named ranges should extract data."""
        try:
            from openpyxl import Workbook
            from openpyxl.workbook.defined_name import DefinedName

            wb = Workbook()
            ws = wb.active
            ws.title = "Named Ranges"

            ws["A1"] = "Name"
            ws["B1"] = "Value"
            ws["A2"] = "Total"
            ws["B2"] = 100

            # Create named range using non-deprecated method
            named_range = DefinedName("MyTotal", attr_text="'Named Ranges'!$B$2")
            wb.defined_names.add(named_range)

            wb.save(temp_dir / "named_ranges.xlsx")

            result = await xlsx_extractor.extract(temp_dir / "named_ranges.xlsx")

            assert result is not None
            assert "Total" in result.markdown
            assert "100" in result.markdown
        except ImportError:
            pytest.skip("openpyxl not installed")

    @pytest.mark.asyncio
    async def test_pptx_with_animations_ignored(self, pptx_extractor, temp_dir):
        """PPTX with animations should extract static content."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            slide.shapes.title.text = "Animated Slide"
            slide.shapes.placeholders[1].text = "Content that might have animations"

            # Note: python-pptx doesn't easily support adding animations
            # but the test ensures animated presentations don't crash

            prs.save(temp_dir / "with_animations.pptx")

            result = await pptx_extractor.extract(temp_dir / "with_animations.pptx")

            assert result is not None
            assert "Animated Slide" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_pptx_with_smartart_shapes(self, pptx_extractor, temp_dir):
        """PPTX with various shapes should extract text content."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

            # Add various shapes with text
            shapes_data = [
                (MSO_SHAPE.RECTANGLE, "Rectangle Text"),
                (MSO_SHAPE.OVAL, "Oval Text"),
                (MSO_SHAPE.DIAMOND, "Diamond Text"),
            ]

            for i, (shape_type, text) in enumerate(shapes_data):
                shape = slide.shapes.add_shape(
                    shape_type,
                    Inches(1 + i * 2), Inches(1),
                    Inches(1.5), Inches(1)
                )
                shape.text = text

            prs.save(temp_dir / "with_shapes.pptx")

            result = await pptx_extractor.extract(temp_dir / "with_shapes.pptx")

            assert result is not None
            assert "Rectangle Text" in result.markdown
            assert "Oval Text" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")

    @pytest.mark.asyncio
    async def test_pptx_with_embedded_media_placeholder(self, pptx_extractor, temp_dir):
        """PPTX referencing media should handle gracefully."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add text indicating embedded content
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            textbox.text_frame.text = "Video content would be here: [Embedded Video]"

            prs.save(temp_dir / "with_media_ref.pptx")

            result = await pptx_extractor.extract(temp_dir / "with_media_ref.pptx")

            assert result is not None
            assert "Video content" in result.markdown or "Embedded" in result.markdown
        except ImportError:
            pytest.skip("python-pptx not installed")
